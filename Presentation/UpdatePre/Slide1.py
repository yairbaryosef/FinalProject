from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

class Slide1Editor:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        self.slide = self.presentation.slides[0]
        self.title_style = {}
        self.content_style = {}

    def _extract_styles(self):
        for shape in self.slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if "Market Analysis" in text and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    self.title_style = {
                        "size": run.font.size,
                        "bold": run.font.bold,
                        "color": run.font.color.rgb
                    }
                elif "Welcome to our" in text and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    self.content_style = {
                        "size": run.font.size,
                        "bold": run.font.bold,
                        "color": run.font.color.rgb
                    }

    def update_title_and_content(self, new_title, new_content):
        self._extract_styles()

        for shape in self.slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                original_text = shape.text.strip()
                tf.clear()

                if "Market Analysis" in original_text or "Hey Yair" in original_text:
                    run = tf.paragraphs[0].add_run()
                    run.text = new_title
                    run.font.size = self.title_style.get("size")
                    run.font.bold = self.title_style.get("bold")
                    run.font.color.rgb = self.title_style.get("color")

                elif "Welcome to our" in original_text or "Hey hey hey" in original_text:
                    run = tf.paragraphs[0].add_run()
                    run.text = new_content
                    run.font.size = self.content_style.get("size")
                    run.font.bold = self.content_style.get("bold")
                    run.font.color.rgb = self.content_style.get("color")

    def save(self, output_path):
        self.presentation.save(output_path)
