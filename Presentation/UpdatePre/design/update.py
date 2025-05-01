# ××¤×¢×™×œ ××ª ×”×§×•×“ ×”××ª×•×§×Ÿ ×¢×œ ×”×§×•×‘×¥ ×©×”××©×ª××© ×”×¢×œ×”

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor


class PresentationBuilder:
    def __init__(self,prs):
        self.prs = prs
        pass

   # ×¤×•× ×§×¦×™×” ×œ×™×¦×™×¨×ª ×©×§×£ ×’×¨×£ ×—×“×©
    def create_column_chart_slide_simple(self,prs, data_dict, max_scale=None):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # ×¨×§×¢ ×œ×‘×Ÿ
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # ×›×•×ª×¨×ª ×‘×× ×¨
        title_shape = slide.shapes.title
        title_shape.text = "Financial Ratios"
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # ×’×¨×£ ×¢××•×“×•×ª (×œ× ××•×¤×§×™)
        chart_data = CategoryChartData()
        chart_data.categories = list(data_dict.keys())
        chart_data.add_series('', list(data_dict.values()))
        print(chart_data.categories.leaf_count)
        x, y, cx, cy = Inches(0.8), Inches(2.2), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

        # ×¢×™×¦×•×‘ ×’×¨×£
        chart.has_legend = False
        if max_scale:
            chart.value_axis.maximum_scale = max_scale

        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.category_axis.tick_labels.font.color.rgb = RGBColor(10, 10, 10)
        chart.value_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.font.color.rgb = RGBColor(10, 10, 10)

        if chart.value_axis.has_major_gridlines:
            gridlines = chart.value_axis.major_gridlines
            gridlines.format.line.color.rgb = RGBColor(200, 200, 200)
            gridlines.format.line.width = Pt(0.75)

        return slide
    def create_CAGR_slide(self,timeseries_data,years):

        # ×—×™×©×•×‘ ×©×™× ×•×™ ×‘××—×•×–×™× ××©× ×” ×§×•×“××ª (2020-2023 => 2021-2023)
        # ×—×™×©×•×‘ CAGR ×›×œ ×©× ×” (×‘×¤×•×¢×œ: ×©×™× ×•×™ ×©× ×ª×™ ××“×•×¨×’)
        def calc_cagr_from_start(values):
            base = values[0]
            return [0.0] + [
                round(((v / base) ** (1 / i) - 1) * 100, 2)
                for i, v in enumerate(values[1:], start=1)
            ]

        pct_change_data = {k: calc_cagr_from_start(v) for k, v in timeseries_data.items()}

         # ×©×™× ×•×™ ×‘××—×•×–×™× ××”×©× ×” ×”×§×•×“××ª

        # ×™×¦×™×¨×ª ×”×©×§×•×¤×™×ª
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # ×›×•×ª×¨×ª ×‘×“×™×•×§ ×›××• ×‘×©×§×•×¤×™×•×ª ×”×›×—×•×œ×•×ª
        title_shape = slide.shapes.title
        title_shape.text = "Financial Change % Over Time"
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # ×™×¦×™×¨×ª × ×ª×•× ×™ ×’×¨×£
        chart_data = CategoryChartData()
        chart_data.categories = years
        for name, values in pct_change_data.items():
            chart_data.add_series(name, values)

        # ×’×¨×£ ××•×¡×˜ ×œ××˜×”
        x, y, cx, cy = Inches(0.6), Inches(2.5), Inches(9), Inches(4.7)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data).chart

        # ×¢×™×¦×•×‘
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(200, 200, 200)

    def create_sentiment_donut_chart(self, sentiment_data: dict):
        from pptx.enum.chart import XL_LEGEND_POSITION

        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # ×›×•×ª×¨×ª
        title_shape = slide.shapes.title
        title_shape.text = "Sentiment Analysis Distribution"
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # ×”×’×“×¨×•×ª ×’×¨×£
        chart_data = CategoryChartData()
        chart_data.categories = list(sentiment_data.keys())
        chart_data.add_series('Sentiment', list(sentiment_data.values()))

        chart_width = Inches(5.5)
        chart_height = Inches(4.0)

        x = Inches(3.5)
        y = Inches(2.5)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, chart_width, chart_height, chart_data
        ).chart

        chart.has_title = False

        # ×¦×‘×¢×™×
        slice_colors = {
            "Positive": RGBColor(0, 176, 80),  # ×™×¨×•×§
            "Neutral": RGBColor(128, 128, 128),  # ××¤×•×¨
            "Negative": RGBColor(255, 0, 0)  # ××“×•×
        }

        # âœ… ×”×•×¡×¤×ª ××—×•×–×™× ×¢×œ ×”×’×¨×£
        series = chart.plots[0].series[0]
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_percentage = True
        data_labels.show_percentage = True
        data_labels.font.size = Pt(14)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)

        # ×¦×‘×¢×™ ×”×¤×œ×— ×œ×¤×™ ×¡×•×’ ×¨×’×©
        for i, category in enumerate(sentiment_data.keys()):
            point = series.points[i]
            if category in slice_colors:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = slice_colors[category]

        # ××’×“×”
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)

    def create_swot_slide(self, title: str, section_data: dict):
        from pptx.enum.text import PP_ALIGN

        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        # ×›×•×ª×¨×ª ×”×©×§×•×¤×™×ª
        slide.shapes.title.text = title

        # ××™×§×•× ×—×“×© â€” × ××•×š ×™×•×ª×¨
        left = Inches(1)
        top = Inches(2.5)  # ğŸ‘ˆ ×ª×™×‘×ª ×”×˜×§×¡×˜ ××•×¡×˜×ª ×œ××˜×”
        width = Inches(8)
        height = Inches(5)

        # ×™×¦×™×¨×ª ×ª×™×‘×ª ×˜×§×¡×˜ ×—×“×©×”
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        for section_title, bullet_points in section_data.items():
            # ×›×•×ª×¨×ª ×”×¡×¢×™×£
            p = text_frame.add_paragraph()
            p.text = section_title
            p.font.bold = True
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)

            # ×¤×¡×§×” ××—×ª ×¢× ×›×œ ×”×‘×•×œ×˜×™×
            bullet_text = "\n".join(f"â€¢ {bp}" for bp in bullet_points)
            para = text_frame.add_paragraph()
            para.text = bullet_text
            para.font.size = Pt(16)
            para.alignment = PP_ALIGN.LEFT
            para.line_spacing = Pt(24)  # ğŸ‘ˆ ×©×•×¨×” ×•×—×¦×™ ×¢×‘×•×¨ ×˜×§×¡×˜ ×‘×’×•×“×œ 16pt

    def create_stock_forecast_slide(self, forecast_data: dict):
        from pptx.enum.chart import XL_LEGEND_POSITION

        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # ×›×•×ª×¨×ª ×”×©×§×•×¤×™×ª
        title_shape = slide.shapes.title
        title_shape.text = "Stock Price Forecast"
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # ×™×¦×™×¨×ª ×’×¨×£ ×¢××•×“×•×ª
        chart_data = CategoryChartData()
        chart_data.categories = list(forecast_data.keys())
        chart_data.add_series("Forecast (%)", list(forecast_data.values()))

        x, y, cx, cy = Inches(1), Inches(2.3), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

        # ×¢×™×¦×•×‘
        chart.has_legend = False
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.font.size = Pt(12)

        # ×¦×‘×¢×™× ×œ×¢××•×“×•×ª: ×™×¨×•×§ ×× ×—×™×•×‘×™, ××“×•× ×× ×©×œ×™×œ×™
        for i, val in enumerate(forecast_data.values()):
            point = chart.plots[0].series[0].points[i]
            point.format.fill.solid()
            if val >= 0:
                point.format.fill.fore_color.rgb = RGBColor(0, 176, 80)  # ×™×¨×•×§
            else:
                point.format.fill.fore_color.rgb = RGBColor(255, 0, 0)  # ××“×•×

    def create_summary_slide(self, recommendation: str, reasons: list):
        from pptx.enum.text import PP_ALIGN

        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # ×›×•×ª×¨×ª ×©×§×•×¤×™×ª
        title_shape = slide.shapes.title
        title_shape.text = "Investment Summary & Recommendation"
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # ××™××•×’â€™×™ ×•×¦×‘×¢
        if recommendation.lower() == "buy":
            rec_text = "Recommendation: BUY ğŸ“ˆ"
            rec_color = RGBColor(0, 176, 80)
        elif recommendation.lower() == "sell":
            rec_text = "Recommendation: SELL ğŸ“‰"
            rec_color = RGBColor(255, 0, 0)
        else:
            rec_text = f"Recommendation: {recommendation}"
            rec_color = RGBColor(0, 0, 0)

        # âœ… ×ª×™×‘×ª ×”××œ×¦×” â€“ ××ª×—×ª ×œ×›×•×ª×¨×ª ××‘×œ ××¢×œ ×”×ª×•×›×Ÿ
        rec_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        tf = rec_box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = rec_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = rec_color
        p.alignment = PP_ALIGN.LEFT

        # âœ… ×ª×™×‘×ª ×‘×•×œ×˜×™× â€“ ××ª×—×™×œ×” × ××•×š ×™×•×ª×¨ ×›×“×™ ×œ× ×œ×”×ª× ×’×©
        reasons_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3.5))
        tf2 = reasons_box.text_frame
        tf2.word_wrap = True
        tf2.clear()

        bullet_text = "\n".join(f"â€¢ {reason}" for reason in reasons)
        para = tf2.add_paragraph()
        para.text = bullet_text
        para.font.size = Pt(16)
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = Pt(30)

    def update_company_overview_slide(self, lines):
        from pptx.enum.text import PP_ALIGN

        for slide in self.prs.slides:
            if slide.shapes.title and "company overview" in slide.shapes.title.text.strip().lower():
                slide.shapes.title.text = "Company Overview"

                # ×”×¡×¨×ª ×ª×™×‘×•×ª ×˜×§×¡×˜ ×§×™×™××•×ª (××œ×‘×“ ×”×›×•×ª×¨×ª)
                shapes_to_remove = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)

                # ×ª×™×‘×ª ×˜×§×¡×˜ ×—×“×©×” â€” ××™×§×•× × ××•×š ×™×•×ª×¨
                textbox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(4.5))
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                text_frame.clear()

                for i, line in enumerate(lines):
                    para = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                    para.text = line
                    para.font.size = Pt(24 if not line.startswith("â€¢") else 22)
                    para.alignment = PP_ALIGN.LEFT
                    para.line_spacing = Pt(36)
                return

        # ×× ×”×©×§×•×¤×™×ª ×œ× ×§×™×™××ª â€” ×¦×•×¨ ×©×§×•×¤×™×ª ×—×“×©×”
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Company Overview"

        textbox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(4.5))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        for i, line in enumerate(lines):
            para = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            para.text = line
            para.font.size = Pt(24 if not line.startswith("â€¢") else 22)
            para.alignment = PP_ALIGN.LEFT
            para.line_spacing = Pt(36)

    def update_cover_slide(self, company_name: str, presenters: str = "YAIR, DAVID AND OREL", date_text: str = None):
        from datetime import datetime
        from pptx.enum.text import PP_ALIGN

        if date_text is None:
            date_text = datetime.today().strftime("%d/%m/%Y")

        slide = self.prs.slides[0]

        # ××—×™×§×ª ×ª×™×‘×•×ª ×˜×§×¡×˜ ×§×™×™××•×ª (×—×•×¥ ××”×›×•×ª×¨×ª)
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                sp = shape._element
                sp.getparent().remove(sp)

        # ğŸ¨ ×¦×‘×¢×™×
        white = RGBColor(255, 255, 255)
        light_blue = RGBColor(102, 153, 255)

        # ×›×•×ª×¨×ª: Stock Financial Analysis (×©×•×¨×” ×¨××©×•× ×”)
        slide.shapes.title.text = f"Stock Financial Analysis:\n{company_name}"
        title_frame = slide.shapes.title.text_frame

        # ×©×•×¨×” 1
        p1 = title_frame.paragraphs[0]
        p1.font.size = Pt(50)
        p1.font.bold = True
        p1.font.color.rgb = white
        p1.alignment = PP_ALIGN.CENTER

        # ×©×•×¨×” 2 (×©× ×”×—×‘×¨×”)
        if len(title_frame.paragraphs) < 2:
            title_frame.add_paragraph()

        p2 = title_frame.paragraphs[1]
        p2.text = company_name
        p2.font.size = Pt(54)
        p2.font.bold = True
        p2.font.color.rgb = white
        p2.alignment = PP_ALIGN.CENTER

        # ×ª×™×‘×ª ×˜×§×¡×˜ ×œ××¦×™×’×™× + ×ª××¨×™×š
        textbox = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(1.2))
        tf = textbox.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.text = f"PRESENTED BY {presenters.upper()}\n{date_text}"
        para.font.size = Pt(24)
        para.font.color.rgb = light_blue
        para.alignment = PP_ALIGN.CENTER


def Create(lines,new_ratios, timeseries_data,years, sentiment, section_data_positive, section_data_negative,
           forecast, recommendation, reasons):
    pptx_path = 'Stock Financial Analysis.pptx'
    prs = Presentation(pptx_path)


    # ×¤×™×œ×•×—
    small_scale_ratios = {k: v for k, v in new_ratios.items() if  v <= 2}
    large_scale_ratios = {k: v for k, v in new_ratios.items() if  v > 2}
    print(small_scale_ratios)
    print(large_scale_ratios)
    # ××—×™×§×ª ×›×œ ×”×©×§×¤×™× ×©×§×©×•×¨×™× ×œ-"Financial Ratios"
    slides_to_remove = []
    for idx, slide in enumerate(prs.slides):
        if slide.shapes.title and "Financial Ratios" in slide.shapes.title.text:
            slides_to_remove.append(idx)

    for idx in sorted(slides_to_remove, reverse=True):
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[idx])
    # ×™×¦×™×¨×” ×‘×¤×•×¢×œ
    builder = PresentationBuilder(prs)
    builder.update_cover_slide("DemoTech Inc.")

    builder.update_company_overview_slide(lines)


    builder.create_column_chart_slide_simple(prs, small_scale_ratios, max_scale=2)
    builder.create_column_chart_slide_simple(prs, large_scale_ratios)
    # ×™×¦×™×¨×ª ×©×§×•×¤×™×ª ×¢× ×’×¨×£ ×§×•×•×™ ×©×œ ×©×™× ×•×™ ×‘××—×•×–×™× ×œ××•×¨×š 4 ×©× ×™×

    # ×¢×¨×›×™× ××§×•×¨×™×™× ×œ×“×•×’××”:

    builder.create_CAGR_slide(timeseries_data,years)

    builder.create_sentiment_donut_chart(sentiment)
    builder.create_swot_slide("Strengths and Opportunities", section_data_positive)

    builder.create_swot_slide("Weaknesses and Threats",section_data_negative )

    builder.create_stock_forecast_slide(forecast)
    builder.create_summary_slide(recommendation, reasons)

    # ×©××™×¨×”
    prs.save("Stock_Financial_Analysis_Final_With_PctChangeChart.pptx")
    print("Process Complete")
