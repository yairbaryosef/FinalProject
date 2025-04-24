from pptx import Presentation
from pptx.dml.color import RGBColor

class Slide2Editor:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        self.slide = self.presentation.slides[1]

    def reset_and_update_content(self,section_texts):

        flag1=False
        flag2=False
        flag3=False
        for shape in self.slide.shapes:
            if not shape.has_text_frame:
                continue

            first_para = shape.text_frame.paragraphs[0].text.lower()
            print(first_para)
            if "founded" in first_para or "apollo" in first_para:
                section = "Digital Pioneer"

            elif "email" in first_para or "content" in first_para:
                section = "Key Business Segments"

            elif "meta" in first_para or "million" in first_para:
                section = "Market Position"
            elif "advertising" in first_para or "media" in first_para:
                shape.text_frame.clear()
                continue
            else:
                continue

            shape.text_frame.clear()
            shape.text_frame.paragraphs[0].text = section
            if (flag1 and section == "Digital Pioneer") or (flag2 and section == "Key Business Segments") or (flag3 and section == "Market Position"):
              continue
            for line in section_texts[section]:
                p = shape.text_frame.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.color.rgb = RGBColor(255, 255, 255)
            if section == "Digital Pioneer":
                flag1 = True
            elif section == "Key Business Segments":
                flag2=True
            elif section == "Market Position":
                flag3=True
    def save(self, output_path):
        self.presentation.save(output_path)
