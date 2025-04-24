import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import yfinance as yf
from pptx.dml.color import RGBColor


from Anlysis.Finnancial_Ratios.scraperyahoo import calc_Ratios_with_growth
from SSH.connect_to_slurm import ssh_connect_and_authenticate


def extract_introduction(summary_path):
    with open(summary_path, "r", encoding="utf-8") as file:
        text = file.read()
    split_text = "Write the introduction in a clear, professional, and concise manner, as if presenting the company"
    parts = text.split(split_text)
    return parts[1].strip() if len(parts) > 1 else ""


def load_sentiment_data(sentiment_path):
    with open(sentiment_path, "r", encoding="utf-8") as file:
        return json.load(file)


def classify_sentences(sentences, criteria):
    # Initialize classification dictionary with empty sets for deduplication
    classification = {criterion: {"positive": set(), "negative": set(), "neutral": set()} for criterion in criteria}

    for entry in sentences:
        sentence = entry["sentence"]
        sentiment = entry.get("predicted_sentiment", "neutral").lower()

        # Classify the sentence to get its best matching category
        sentence_category_mapping, _ = clasify([sentence])
        best_match_category = list(sentence_category_mapping.keys())[0] if sentence_category_mapping else None

        # Only keep the sentence if it matched a valid financial criterion
        if best_match_category in classification:
            classification[best_match_category][sentiment].add(sentence)

    # Convert sets back to lists before returning
    for category in classification:
        for sentiment in classification[category]:
            classification[category][sentiment] = list(classification[category][sentiment])

    return classification

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
# Function to classify documents (sentences in this case) and map them to the best category based on cosine similarity
def clasify(documents):
    sentiment_criteria = [
        "Non-Financial","Order Backlog", "Revenue Growth", "Net Income", "Operating Income", "Gross Profit",
        "Product Demand", "Awarded Contracts", "Market Expansion", "Revenue Diversification", "Market Share",
        "Research and Development", "Trade Payables", "Contract Liabilities", "Strategic Shifts",
        "Restructuring Losses", "Cash Flow Decrease", "Operating Expenses", "Intersegment Revenue",
        "Trade and Unbilled Receivables", "Supply Chain Disruptions", "Geopolitical Impact",
        "Production Relocation", "Stock-Based Compensation", "Forward-Looking Statements", "Financial Reporting"
    ]

    # Vectorizing the criteria and documents
    vectorizer = TfidfVectorizer()
    vectorized_texts = vectorizer.fit_transform(sentiment_criteria + documents)

    # Separate category vectors and document vectors
    category_vectors = vectorized_texts[:len(sentiment_criteria)]
    document_vectors = vectorized_texts[len(sentiment_criteria):]
    cosine_similarities = cosine_similarity(document_vectors, category_vectors)

    # Mapping documents to the best matching category based on cosine similarity
    doc_category_mapping = {}
    for i, doc in enumerate(documents):
        best_match_idx = np.argmax(cosine_similarities[i])
        best_category = sentiment_criteria[best_match_idx]
        if best_category not in doc_category_mapping:
            doc_category_mapping[best_category] = []
        doc_category_mapping[best_category].append(doc)

    return doc_category_mapping, sentiment_criteria


from pptx.util import Inches

from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def add_slide(prs, title, sentences_by_sentiment, max_lines=10):
    """Adds slides with categorized sentiment-based content and a gradient background."""

    sentiment_colors = {
        "positive": RGBColor(0, 128, 0),  # Green
        "negative": RGBColor(255, 0, 0),  # Red
        "neutral": RGBColor(0, 0, 0)  # Black (default)
    }


    for sentiment, sentences in sentences_by_sentiment.items():
        if not sentences:
            continue  # Skip empty categories

        total_sentences = len(sentences)
        num_slides = (total_sentences // max_lines) + (1 if total_sentences % max_lines else 0)

        for slide_idx in range(num_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout
            if sentiment == "positive":
                gradient_colors = [(0, 255, 0), (34, 139, 34)]  # Green gradient
            elif sentiment == "negative":
                gradient_colors = [(255, 0, 0), (139, 0, 0)]  # Red gradient
            else:
                gradient_colors = [(169, 169, 169), (211, 211, 211)]  # Grey gradient

            # Apply gradient background
            background = slide.background
            fill = background.fill
            fill.gradient()  # Enable gradient
            fill.gradient_stops[0].position = 0.0
            fill.gradient_stops[0].color.rgb = RGBColor(*gradient_colors[0])  # Dark blue
            fill.gradient_stops[1].position = 1.0
            fill.gradient_stops[1].color.rgb = RGBColor(*gradient_colors[1])  # Lighter blue

            # Title
            title_shape = slide.shapes.title
            title_shape.text = f"{title} - {sentiment.capitalize()} ({slide_idx + 1}/{num_slides})"
            title_shape.text_frame.paragraphs[0].font.size = Pt(20)
            # Content box
            content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(50), Pt(100), Pt(600), Pt(350))
            content_box.fill.solid()
            content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for readability
            content_box.shadow.inherit = False

            text_frame = content_box.text_frame
            text_frame.clear()  # Ensure it's empty before adding new content

            # Add sentiment category header
            p = text_frame.add_paragraph()
            p.text = f"{sentiment.capitalize()} Sentences:"
            p.font.bold = True
            p.font.size = Pt(14)

            # Add sentences with respective colors
            start_idx = slide_idx * max_lines
            end_idx = min(start_idx + max_lines, total_sentences)

            for sentence in sentences[start_idx:end_idx]:
                p = text_frame.add_paragraph()
                p.text = f"- {sentence}"
                p.font.size = Pt(14)
                p.font.color.rgb = sentiment_colors.get(sentiment, RGBColor(0, 0, 0))  # Default to black
            if "Order Backlog" in title :
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "backlog.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Revenue Growth" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "Revenue.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Net Income" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "NetIncome.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Operating Income" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "factory.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Gross Profit" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "gross.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Awarded Contracts" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "contract.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Research and Development" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "RD.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Contract Liabilities" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "liability-insurance.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Operating Expenses" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "opearting_expenses.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Intersegment Revenue" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "self_invest.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Trade and Unbilled Receivables" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "Trade and Unbilled Receivables.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Supply Chain Disruptions" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "supply-chain.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Stock-Based Compensation" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "stock.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)
            elif "Forward-Looking Statements" in title:
                base_dir = Path("").resolve()
                image_path = base_dir / "Files" / "Pictures" / "goal.png"
                img_width = Inches(1.25)  # Reduced width (half the original size)
                img_height = Inches(1)  # Reduced height (half the original size)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                img_left = slide_width - img_width - Inches(0.5)  # Right bottom with margin
                img_top = slide_height - img_height - Inches(0.5)  # Bottom margin

                slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)




def add_text_slide(prs, title, content):
    """Adds a slide with a single block of text (e.g., Company Introduction)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(18)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    text_box = slide.shapes.add_textbox(Pt(50), Pt(100), Pt(600), Pt(400))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(14)

def create_presentation(summary_path, sentiment_path, output_path,ratios,symbol):
        prs = Presentation()

        # Extract company introduction
        introduction = extract_introduction(summary_path)

        # Add company title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = symbol
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)

        # Add introduction slide
        add_text_slide(prs, "Company Introduction", introduction)

        # Load sentiment data and classify sentences
        sentiment_data = load_sentiment_data(sentiment_path)

        sentiment_criteria = [
            "Non-Financial","Order Backlog", "Revenue Growth", "Net Income", "Operating Income", "Gross Profit",
            "Product Demand", "Awarded Contracts", "Market Expansion", "Revenue Diversification", "Market Share",
            "Research and Development", "Trade Payables", "Contract Liabilities", "Strategic Shifts",
            "Restructuring Losses", "Cash Flow Decrease", "Operating Expenses", "Intersegment Revenue",
            "Trade and Unbilled Receivables", "Supply Chain Disruptions", "Geopolitical Impact",
            "Production Relocation", "Stock-Based Compensation", "Forward-Looking Statements", "Financial Reporting"
        ]

        print(len(sentiment_data))

        # Add a "Non-Financial" category for unclassified sentences
        classified_sentences = classify_sentences(sentiment_data, sentiment_criteria)

        # Remove the "Non-Financial" category after classification
        if "Non-Financial" in classified_sentences:
            del classified_sentences["Non-Financial"]

        # Add slides for financial sentiment analysis
        for category, sentiments in classified_sentences.items():
            add_slide(prs, category, sentiments)

        # Save presentation

        # Financial Ratios Slide (Placeholder)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Financial Ratios"
        content = slide.placeholders[1]
        text = ""
        for ratio, value in ratios.items():
            try:
                text += f"{ratio}: {value:.2f}\n"
            except:
                text += f"{ratio}: None\n"
        content.text = text
        for p in content.text_frame.paragraphs:
         p.font.size = Pt(12)
        # Save Presentation
        prs.save(output_path)

import os
def main(stock,symbol,pdf):


    ssh_connect_and_authenticate(local_file_path=pdf)
    try:
     ratios=calc_Ratios_with_growth(stock,symbol)
    except:
        ratios={}

    base_dir = os.path.dirname(os.path.abspath(""))

    summary_file = os.path.join(base_dir, "Presentation", "Files", "company_summary.txt")

    # Construct the full path for the sentiment_results.json file inside Presentation/Files
    sentiment_file = os.path.join(base_dir, "Presentation", "Files", "sentiment_results.json")

    output_pptx = os.path.join(base_dir, "Presentation", "Files","Company_Presentation.pptx")
    print(output_pptx)
    # Run the script
    create_presentation(summary_file, sentiment_file, output_pptx,ratios,symbol)
    return output_pptx



def process(symbol,local_file_path):

    print(local_file_path)
    stock = yf.Ticker(symbol)
    return main(stock,symbol,local_file_path)
    # Create the presentation
    #return create_presentation(ratios,symbol,stock,local_file_path)
if __name__ == '__main__':
    process('ESLT.TA',"")
